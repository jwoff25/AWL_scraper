import pptx
import re

ppt_list = ['sub1.pptx', 'sub2.pptx', 'sub3.pptx', 'sub4.pptx']
col_list = []
counter = 1;

col_regex = re.compile('Collocation[\s]*:')
space_regex = re.compile('\s{2,}')

with open("collocation.txt", 'a') as coltext:
	for sub in ppt_list:
		p = pptx.Presentation(sub)
		slides = p.slides
		for sl in slides:
			shape_list = sl.shapes
			text_obj = [s.text_frame for s in shape_list if s.has_text_frame]
			text_list = [t.text.encode('ascii', 'ignore').strip() for t in text_obj]
			text_list = [t.replace("\n", ";") for t in text_list]
			split_list = [t.split(";") for t in text_list]
			collocation = [c for li in split_list for c in li if re.match('^Collocation.', c)]
			collocation = [re.sub(col_regex, "", c) for c in collocation]
			collocation = [re.sub(space_regex, " ", c) for c in collocation]
			for col in collocation:
				col_list.append("".join(col))
				coltext.write('"' + "".join(col) + '"\n')
			print "----------------"

print col_list
